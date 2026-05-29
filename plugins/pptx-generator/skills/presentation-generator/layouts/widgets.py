"""widgets — composable slide elements.

Reusable building blocks for complex layouts. Each function takes a slide,
a position (left/top/width/height as fractions of slide W/H), and content
parameters, and adds shapes/text to the slide.

Conventions:
  * All positions and sizes are fractions (0.0–1.0) of slide dimensions.
  * Functions are pure: they add shapes but don't read them back.
  * Colours come from `colors.py` helpers — never hardcoded.
"""
from __future__ import annotations

from typing import Optional, Iterable

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from . import colors


def _emu_xy(prs, l, t, w, h):
    W, H = prs.slide_width, prs.slide_height
    return Emu(int(W * l)), Emu(int(H * t)), Emu(int(W * w)), Emu(int(H * h))


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#"))


# Default font applied to every text run unless overridden. Set once from
# brand.yaml via `set_default_fonts()` at the start of generation so every
# widget renders in Raleway / Raleway SemiBold instead of falling back to
# Arial.
_DEFAULT_FONTS = {
    "body":     "Raleway",
    "emphasis": "Raleway SemiBold",
    "light":    "Raleway Light",
    "fallback": "Arial",
}


def set_default_fonts(*, body=None, emphasis=None, light=None, fallback=None):
    if body:     _DEFAULT_FONTS["body"]     = body
    if emphasis: _DEFAULT_FONTS["emphasis"] = emphasis
    if light:    _DEFAULT_FONTS["light"]    = light
    if fallback: _DEFAULT_FONTS["fallback"] = fallback


def _set_text(tf, lines, size_pt, bold=False, color="#000000", align="left",
              font=None):
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    font_name = font or (_DEFAULT_FONTS["emphasis"] if bold else _DEFAULT_FONTS["body"])
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(line)
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        for run in p.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.name = font_name
            run.font.color.rgb = _rgb(color)


def text_box(slide, prs, l, t, w, h, text, *,
             size_pt=14, bold=False, color="#000000",
             align="left", anchor="top", wrap=True, font=None):
    """Plain text box; the workhorse."""
    left, top, width, height = _emu_xy(prs, l, t, w, h)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(int(prs.slide_width * 0.003))
    tf.margin_top = tf.margin_bottom = Emu(int(prs.slide_height * 0.003))
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    lines = str(text).split("\n")
    _set_text(tf, lines, size_pt, bold=bold, color=color, align=align, font=font)
    return box


def filled_rect(slide, prs, l, t, w, h, *, fill_color, line_color=None):
    """Filled rectangle (used as background, divider, accent bar)."""
    left, top, width, height = _emu_xy(prs, l, t, w, h)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    if line_color:
        shape.line.color.rgb = _rgb(line_color)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def rounded_rect(slide, prs, l, t, w, h, *, fill_color, line_color=None):
    left, top, width, height = _emu_xy(prs, l, t, w, h)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    if line_color:
        shape.line.color.rgb = _rgb(line_color)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def circle(slide, prs, l, t, size, *, fill_color, line_color=None):
    left, top, width, height = _emu_xy(prs, l, t, size, size)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    if line_color:
        shape.line.color.rgb = _rgb(line_color)
    else:
        shape.line.fill.background()
    return shape


# ── Composite widgets ─────────────────────────────────────────────────────

def kpi_tile(slide, prs, l, t, w, h, *,
             value, label, brand,
             trend_delta=None, status_key=None):
    """Big number with caption; optional trend arrow and status colour.

    Layout:
        ┌───────────────┐
        │   VALUE  ↑    │   ← large number + optional arrow
        │   label       │   ← caption
        └───────────────┘
    """
    color = colors.status(status_key, default=brand["colors"]["primary"]) if status_key else brand["colors"]["primary"]

    # Background tile
    rounded_rect(slide, prs, l, t, w, h,
                 fill_color="#FFFFFF",
                 line_color=brand["colors"]["gray_2"])

    arrow_str = ""
    if trend_delta is not None:
        arrow, trend_color = colors.trend(trend_delta)
        arrow_str = f" {arrow}"
        color = trend_color  # let trend override status

    text_box(slide, prs,
             l + 0.005, t + 0.005, w - 0.01, h * 0.55,
             text=f"{value}{arrow_str}",
             size_pt=28, bold=True, color=color, align="left", anchor="middle")
    text_box(slide, prs,
             l + 0.005, t + h * 0.55, w - 0.01, h * 0.40,
             text=label,
             size_pt=11, color=brand["colors"]["gray_1"], align="left", anchor="top")


def person_card(slide, prs, l, t, w, h, *,
                name, role, bio, brand, accent_color=None):
    """Team member tile: name (bold) + role (red) + 1-sentence bio."""
    accent = accent_color or brand["colors"]["primary"]
    text_box(slide, prs, l, t, w, h * 0.20,
             text=name, size_pt=15, bold=True,
             color=brand["colors"]["black"], anchor="middle")
    text_box(slide, prs, l, t + h * 0.20, w, h * 0.15,
             text=role, size_pt=11, bold=True,
             color=accent, anchor="middle")
    text_box(slide, prs, l, t + h * 0.38, w, h * 0.60,
             text=bio, size_pt=10,
             color=brand["colors"]["gray_1"], anchor="top")


def big_stat(slide, prs, l, t, w, h, *,
             value, caption, brand, color=None, align="center"):
    """A single huge stat with a caption below."""
    val_color = color or brand["colors"]["primary"]
    text_box(slide, prs, l, t, w, h * 0.65,
             text=value, size_pt=52, bold=True,
             color=val_color, align=align, anchor="middle")
    text_box(slide, prs, l, t + h * 0.65, w, h * 0.30,
             text=caption, size_pt=12,
             color=brand["colors"]["gray_1"], align=align, anchor="top")


def badge(slide, prs, l, t, w, h, *, text, color, text_color="#FFFFFF"):
    """Small filled rounded label."""
    rounded_rect(slide, prs, l, t, w, h, fill_color=color)
    text_box(slide, prs, l, t, w, h,
             text=text, size_pt=10, bold=True,
             color=text_color, align="center", anchor="middle")


def status_pill(slide, prs, l, t, w, h, *, status_key, brand, label=None):
    """Coloured pill showing status (DONE/BLOCKED/etc)."""
    color = colors.status(status_key)
    badge(slide, prs, l, t, w, h,
          text=(label or status_key).upper(),
          color=color, text_color="#FFFFFF")


def timeline_event(slide, prs, l, t, w, h, *,
                   date, label, brand, status_key=None):
    """One event on a timeline: dot + date + label."""
    dot_size = h * 0.20
    dot_color = colors.status(status_key, default=brand["colors"]["primary"])
    circle(slide, prs, l, t, dot_size, fill_color=dot_color)
    text_box(slide, prs, l, t + dot_size + 0.005,
             w, h * 0.30,
             text=date, size_pt=10, bold=True,
             color=brand["colors"]["gray_1"], align="left")
    text_box(slide, prs, l, t + dot_size + 0.04,
             w, h * 0.50,
             text=label, size_pt=12,
             color=brand["colors"]["black"], align="left")


def comparison_row(slide, prs, l, t, w, row_h, *,
                   label, options, brand, label_w=0.20):
    """Row in a comparison matrix: label | option1 | option2 | ...

    options: iterable of booleans (✓/✗) or strings.
    """
    text_box(slide, prs, l, t, label_w, row_h,
             text=label, size_pt=11, bold=True,
             color=brand["colors"]["black"], anchor="middle")
    rest_w = w - label_w
    n = max(1, len(options))
    col_w = rest_w / n
    for i, opt in enumerate(options):
        cx = l + label_w + i * col_w
        if isinstance(opt, bool):
            sym = "✔" if opt else "✘"
            col = brand["colors"]["green"] if opt else brand["colors"]["primary"]
            text_box(slide, prs, cx, t, col_w, row_h,
                     text=sym, size_pt=18, bold=True, color=col,
                     align="center", anchor="middle")
        else:
            text_box(slide, prs, cx, t, col_w, row_h,
                     text=str(opt), size_pt=11,
                     color=brand["colors"]["black"],
                     align="center", anchor="middle")


def section_label(slide, prs, l, t, w, h, *, text, brand):
    """Small caps label, e.g. "WYTYCZNE" / "SOLUTION" used as section eyebrow."""
    text_box(slide, prs, l, t, w, h,
             text=text.upper(), size_pt=10, bold=True,
             color=brand["colors"]["primary"], align="left", anchor="middle")


def divider(slide, prs, l, t, w, *, color, thickness=0.002):
    """Thin horizontal divider line."""
    filled_rect(slide, prs, l, t, w, thickness, fill_color=color)
