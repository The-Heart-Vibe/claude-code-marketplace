"""free_canvas — multi-column / grid layouts drawn as free text boxes.

Some templates (notably blank.pptx) don't ship native multi-column master
layouts. Instead of forcing the spec into a misaligned native layout, we
add a slide using the bare TITLE/BLANK master and place text boxes at
computed positions on it. The result reads visually identical to a native
column layout but with predictable geometry across templates.

The dispatcher chooses between free-canvas and native rendering based on
the template_map cfg. If `cfg['free_canvas']` is set, we use this module;
otherwise the regular generic layouts run.
"""
from __future__ import annotations

from typing import Callable

from pptx.util import Emu

from ..base import add_slide, placeholder, set_text, set_bullets, set_notes


# Margins as fraction of slide dimensions
MARGIN_X = 0.04
MARGIN_Y = 0.07
TITLE_HEIGHT = 0.10
BODY_TOP = 0.20      # where the body grid begins (below title)
BODY_BOTTOM = 0.92   # where the body grid ends (above footer area)
GUTTER = 0.015       # space between columns


def _add_text_box(slide, prs, left_pct, top_pct, w_pct, h_pct, text,
                  *, bold=False, size_pt=14, color=None, align="left"):
    """Add a free-floating text box at percentage coordinates."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    W, H = prs.slide_width, prs.slide_height
    left   = Emu(int(W * left_pct))
    top    = Emu(int(H * top_pct))
    width  = Emu(int(W * w_pct))
    height = Emu(int(H * h_pct))

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    tf.vertical_anchor = MSO_ANCHOR.TOP

    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        for run in p.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            if color:
                run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
    return box


def _add_title(slide, prs, text):
    title_ph = next(
        (ph for ph in slide.placeholders
         if ph.has_text_frame and ph.placeholder_format.type == 1),  # TITLE = 1
        None,
    )
    if title_ph is not None:
        set_text(title_ph, text)
    else:
        _add_text_box(
            slide, prs,
            MARGIN_X, MARGIN_Y,
            1 - 2 * MARGIN_X, TITLE_HEIGHT,
            text, bold=True, size_pt=24,
        )


# ── Multi-column builders ──────────────────────────────────────────────────

def two_column(prs, layout_cfg, spec, brand, ctx):
    slide = add_slide(prs, layout_cfg)
    _add_title(slide, prs, spec.get("title", ""))

    col_w = (1 - 2 * MARGIN_X - GUTTER) / 2
    col_h = BODY_BOTTOM - BODY_TOP

    for i, (side, key) in enumerate((("left", "left"), ("right", "right"))):
        col_spec = spec.get(key, {}) or {}
        heading = col_spec.get("heading", "")
        bullets = col_spec.get("bullets", [])
        left = MARGIN_X + i * (col_w + GUTTER)

        if heading:
            _add_text_box(
                slide, prs,
                left, BODY_TOP,
                col_w, 0.08,
                heading, bold=True, size_pt=18,
                color=brand["colors"].get("primary"),
            )
        body_text = "\n".join(f"• {b}" for b in bullets)
        if body_text:
            _add_text_box(
                slide, prs,
                left, BODY_TOP + 0.09,
                col_w, col_h - 0.09,
                body_text, size_pt=14,
            )

    set_notes(slide, spec.get("notes"))
    return slide


def three_column(prs, layout_cfg, spec, brand, ctx):
    slide = add_slide(prs, layout_cfg)
    _add_title(slide, prs, spec.get("title", ""))

    subtitle = spec.get("subtitle", "")
    subtitle_h = 0.06 if subtitle else 0
    if subtitle:
        _add_text_box(
            slide, prs,
            MARGIN_X, BODY_TOP,
            1 - 2 * MARGIN_X, subtitle_h,
            subtitle, size_pt=14, color=brand["colors"].get("gray_1", "#969696"),
        )

    cols = _normalize_columns(spec)[:3]
    if not cols:
        return slide

    col_w = (1 - 2 * MARGIN_X - 2 * GUTTER) / 3
    grid_top = BODY_TOP + subtitle_h + 0.02
    grid_h = BODY_BOTTOM - grid_top

    for i, col in enumerate(cols):
        left = MARGIN_X + i * (col_w + GUTTER)
        heading = col.get("heading", "")
        body = col.get("body", "")
        if heading:
            _add_text_box(
                slide, prs,
                left, grid_top,
                col_w, 0.15,
                heading, bold=True, size_pt=36,
                color=brand["colors"].get("primary"),
                align="left",
            )
        if body:
            _add_text_box(
                slide, prs,
                left, grid_top + 0.18,
                col_w, grid_h - 0.18,
                body, size_pt=14,
            )

    set_notes(slide, spec.get("notes"))
    return slide


def stat_pair(prs, layout_cfg, spec, brand, ctx):
    slide = add_slide(prs, layout_cfg)
    _add_title(slide, prs, spec.get("title", ""))

    main = spec.get("main_statement", "")
    supporting = spec.get("supporting", [])
    aux = spec.get("aux", "")

    if main:
        _add_text_box(
            slide, prs,
            MARGIN_X, BODY_TOP,
            1 - 2 * MARGIN_X, 0.12,
            main, bold=True, size_pt=20,
        )

    if supporting:
        cols = supporting[:3]
        n = len(cols)
        col_w = (1 - 2 * MARGIN_X - (n - 1) * GUTTER) / n
        grid_top = BODY_TOP + (0.14 if main else 0)
        for i, item in enumerate(cols):
            left = MARGIN_X + i * (col_w + GUTTER)
            label = item.get("label", "")
            body = item.get("body", "")
            if label:
                _add_text_box(
                    slide, prs,
                    left, grid_top,
                    col_w, 0.12,
                    label, bold=True, size_pt=22,
                    color=brand["colors"].get("primary"),
                )
            if body:
                _add_text_box(
                    slide, prs,
                    left, grid_top + 0.14,
                    col_w, BODY_BOTTOM - grid_top - 0.14,
                    body, size_pt=13,
                )

    if aux:
        _add_text_box(
            slide, prs,
            MARGIN_X, BODY_BOTTOM - 0.04,
            1 - 2 * MARGIN_X, 0.04,
            aux, size_pt=9, color=brand["colors"].get("gray_1", "#969696"),
        )

    set_notes(slide, spec.get("notes"))
    return slide


def table_grid(prs, layout_cfg, spec, brand, ctx):
    """Render `cells` as a flexible NxM grid; rows/cols chosen by item count."""
    slide = add_slide(prs, layout_cfg)
    _add_title(slide, prs, spec.get("title", ""))

    cells = list(spec.get("cells", []))
    if not cells:
        set_notes(slide, spec.get("notes"))
        return slide

    n = len(cells)
    cols = 4 if n >= 8 else 3 if n >= 5 else 2 if n >= 3 else n
    rows = (n + cols - 1) // cols

    col_w = (1 - 2 * MARGIN_X - (cols - 1) * GUTTER) / cols
    row_h = (BODY_BOTTOM - BODY_TOP - (rows - 1) * GUTTER) / rows

    for i, cell in enumerate(cells):
        r, c = divmod(i, cols)
        left = MARGIN_X + c * (col_w + GUTTER)
        top = BODY_TOP + r * (row_h + GUTTER)
        _add_text_box(
            slide, prs, left, top, col_w, row_h,
            cell, size_pt=13,
        )

    set_notes(slide, spec.get("notes"))
    return slide


# ── Helpers ────────────────────────────────────────────────────────────────

def _normalize_columns(spec: dict) -> list:
    if "columns" in spec:
        return list(spec["columns"])
    cols = []
    for key in ("col1", "col2", "col3"):
        val = spec.get(key)
        if val is None:
            continue
        if isinstance(val, dict):
            cols.append(val)
        else:
            text = str(val)
            head, _, body = text.partition("\n")
            cols.append({"heading": head.strip(), "body": body.strip()})
    return cols
