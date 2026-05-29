"""content — default slide: title + body (bullets or paragraph).

Spec:
  {
    "layout": "content",
    "title":  "...",
    "body":   ["bullet 1", "bullet 2", ...]   # or a single string
    "icon":   "shield"                         # optional, lucide icon name
  }
"""
from pptx.util import Emu

from ..base import add_slide, placeholder, set_text, set_bullets, set_notes


def build(prs, layout_cfg, spec, brand, ctx):
    slide = add_slide(prs, layout_cfg)

    set_text(placeholder(slide, layout_cfg, "title"), spec.get("title", ""))

    body = spec.get("body", [])
    body_shape = placeholder(slide, layout_cfg, "body")
    if isinstance(body, list):
        set_bullets(body_shape, body)
    elif isinstance(body, str):
        set_text(body_shape, body)

    _inject_icon(slide, spec, brand, prs, ctx)

    set_notes(slide, spec.get("notes"))
    return slide


def _inject_icon(slide, spec, brand, prs, ctx) -> None:
    icon = spec.get("icon")
    if not icon:
        return
    get_icon_png = ctx.get("get_icon_png")
    if not get_icon_png:
        return
    color = brand["colors"]["primary"]
    size = brand["icons"]["default_size"]
    path = get_icon_png(icon, size=size, color=color)
    if not path:
        return
    W, H = prs.slide_width, prs.slide_height
    title = next(
        (sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text),
        None,
    )
    if not title:
        return
    icon_emu = Emu(int(W * 0.055))
    left = title.left
    top = max(0, title.top - icon_emu - Emu(int(H * 0.005)))
    slide.shapes.add_picture(str(path), left, top, icon_emu, icon_emu)
