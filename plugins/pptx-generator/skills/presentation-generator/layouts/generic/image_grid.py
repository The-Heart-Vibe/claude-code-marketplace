"""image_grid — 2x2 grid of captioned images (gallery / proof points).

Spec:
  { "layout": "image_grid",
    "title": "...",
    "tiles": [
      { "image": "/path/to/img1.png", "caption": "Onboarding screen"  },
      { "image": "/path/to/img2.png", "caption": "Payment flow"       },
      { "image": "/path/to/img3.png", "caption": "Restaurant dashboard"},
      { "image": "/path/to/img4.png", "caption": "Analytics overview" }
    ]
  }
  # If `image` paths are missing, a coloured placeholder block is drawn so
  # the layout still reads visually.
"""
from pathlib import Path
from pptx.util import Emu

from ..base import add_slide, set_notes
from ..widgets import text_box, rounded_rect


def build(prs, layout_cfg, spec, brand, ctx):
    slide = add_slide(prs, layout_cfg)

    text_box(slide, prs, 0.05, 0.04, 0.90, 0.08,
             text=spec.get("title", ""), size_pt=22, bold=True,
             color=brand["colors"]["black"])

    tiles = spec.get("tiles", [])[:4]
    if not tiles:
        set_notes(slide, spec.get("notes"))
        return slide

    g_left, g_top, g_w, g_h = 0.05, 0.16, 0.90, 0.78
    cell_w = g_w / 2 - 0.01
    cell_h = g_h / 2 - 0.01

    W, H = prs.slide_width, prs.slide_height

    for i, tile in enumerate(tiles):
        r = i // 2
        c = i % 2
        left = g_left + c * (cell_w + 0.02)
        top = g_top + r * (cell_h + 0.02)
        img_h = cell_h * 0.75

        img_path = tile.get("image")
        if img_path and Path(img_path).exists():
            slide.shapes.add_picture(
                str(img_path),
                Emu(int(W * left)), Emu(int(H * top)),
                Emu(int(W * cell_w)), Emu(int(H * img_h)),
            )
        else:
            rounded_rect(slide, prs, left, top, cell_w, img_h,
                         fill_color=brand["colors"]["gray_3"],
                         line_color=brand["colors"]["gray_2"])
            text_box(slide, prs, left, top, cell_w, img_h,
                     text="(image)", size_pt=12,
                     color=brand["colors"]["gray_1"],
                     align="center", anchor="middle")

        text_box(slide, prs, left, top + img_h + 0.005,
                 cell_w, cell_h - img_h - 0.005,
                 text=tile.get("caption", ""), size_pt=11,
                 color=brand["colors"]["black"], align="center")

    set_notes(slide, spec.get("notes"))
    return slide
