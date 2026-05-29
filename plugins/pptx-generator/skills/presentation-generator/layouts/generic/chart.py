"""chart — title + rendered matplotlib chart image.

Spec:
  {
    "layout": "chart",
    "title": "...",
    "chart_type": "bar" | "line" | "pie",
    "data": { "labels": [...], "values": [...] }
  }
"""
from pathlib import Path

from pptx.util import Emu

from ..base import add_slide, placeholder, set_text, set_notes


def build(prs, layout_cfg, spec, brand, ctx):
    slide = add_slide(prs, layout_cfg)
    set_text(placeholder(slide, layout_cfg, "title"), spec.get("title", ""))

    render_chart = ctx.get("render_chart")
    output_dir   = ctx.get("output_dir")
    if render_chart and output_dir and spec.get("data"):
        path = render_chart(
            spec.get("chart_type", "bar"),
            spec["data"],
            brand,
            str(Path(output_dir) / f"chart_{id(spec)}.png"),
        )
        if path:
            W, H = prs.slide_width, prs.slide_height
            title_shape = placeholder(slide, layout_cfg, "title")
            if title_shape is not None:
                top = title_shape.top + title_shape.height + Emu(int(H * 0.015))
            else:
                top = Emu(int(H * 0.22))
            slide.shapes.add_picture(
                path,
                Emu(int(W * 0.05)),
                top,
                Emu(int(W * 0.90)),
                Emu(int(H * 0.62)),
            )

    set_notes(slide, spec.get("notes"))
    return slide
