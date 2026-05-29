"""brand_chrome — rich programmatic brand visuals applied to every slide.

When templates are stripped to Slide Master only, the brand visuals that
lived as free shapes on the original example slides are lost. This module
re-applies them with multiple variants per slide type so output decks
remain on-brand and visually rich.

Modes (selected via `mode_for(spec)` or explicitly):
  cover_minimal      — small logo + bottom-right triangle, otherwise clean
  cover_split        — red sidebar 25% + logo top of sidebar
  cover_statement    — full red triangle bottom-right + logo top-left
  section            — vertical red bar full height + icon corner
  section_filled     — full coloured background + white logo
  content            — bottom red bar + page number + small horizontal logo
  content_minimal    — thin red top accent + page number, no logo (dense slides)
  closing            — thick red bar + larger logo bottom-right + accent triangle
  data_heavy         — minimalist chrome so dense charts/tables aren't fought

The mode is auto-selected from spec; specs can force a mode via
`"chrome": "cover_statement"` etc.
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


BRAND_ASSETS = Path(__file__).resolve().parent.parent / "brand_assets"
LOGO_H = BRAND_ASSETS / "th-horizontal.png"
LOGO_V = BRAND_ASSETS / "th-vertical.png"
LOGO_ICON = BRAND_ASSETS / "th-icon.png"

# Real PNG aspect ratios (width / height) measured from the source files.
# Using these prevents the logo from being stretched.
LOGO_H_RATIO = 7088 / 2482     # 2.856 — horizontal logo is ~2.86× wider than tall
LOGO_V_RATIO = 2363 / 1900     # 1.244 — "vertical" logo is slightly wider than tall
LOGO_ICON_RATIO = 905 / 1371   # 0.660 — icon is taller than wide


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#"))


# ── Public API ────────────────────────────────────────────────────────────

def apply(slide, prs, brand: dict, *,
          mode: str = "content",
          page_number: Optional[int] = None,
          total_pages: Optional[int] = None,
          section_label: Optional[str] = None) -> None:
    """Apply brand chrome based on the requested mode."""
    fn = MODES.get(mode, _content)
    fn(slide, prs, brand,
       page_number=page_number,
       total_pages=total_pages,
       section_label=section_label)


# ── Cover variants ────────────────────────────────────────────────────────

def _cover_minimal(slide, prs, brand, **kw):
    W, H = prs.slide_width, prs.slide_height
    red = brand["colors"]["primary"]

    # Small accent triangle bottom-right
    tri = Emu(int(W * 0.10))
    t = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, W - tri, H - tri, tri, tri,
    )
    t.fill.solid(); t.fill.fore_color.rgb = _rgb(red); t.line.fill.background()

    # Small horizontal logo top-right
    if LOGO_H.exists():
        logo_h = Emu(int(H * 0.05))
        logo_w = Emu(int(H * 0.05 * LOGO_H_RATIO))
        slide.shapes.add_picture(
            str(LOGO_H),
            W - logo_w - Emu(int(W * 0.04)),
            Emu(int(H * 0.04)),
            logo_w, logo_h,
        )


def _cover_split(slide, prs, brand, **kw):
    W, H = prs.slide_width, prs.slide_height
    red = brand["colors"]["primary"]

    # Red sidebar — left 22% of slide
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(int(W * 0.22)), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(red); bar.line.fill.background()

    # Vertical logo at top of sidebar — use real W/H ratio (1.244)
    if LOGO_V.exists():
        logo_w = Emu(int(W * 0.12))
        logo_h = Emu(int(W * 0.12 / LOGO_V_RATIO))
        slide.shapes.add_picture(
            str(LOGO_V),
            Emu(int(W * 0.05)),
            Emu(int(H * 0.06)),
            logo_w, logo_h,
        )


def _cover_statement(slide, prs, brand, **kw):
    W, H = prs.slide_width, prs.slide_height
    red = brand["colors"]["primary"]

    # Large triangle bottom-right (1/3 of slide)
    tri = Emu(int(W * 0.30))
    t = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, W - tri, H - tri, tri, tri,
    )
    t.fill.solid(); t.fill.fore_color.rgb = _rgb(red); t.line.fill.background()

    # Smaller secondary triangle top-left for balance
    tri2 = Emu(int(W * 0.10))
    t2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, tri2, tri2,
    )
    t2.fill.solid(); t2.fill.fore_color.rgb = _rgb(brand["colors"]["black"])
    t2.line.fill.background()
    t2.rotation = 180

    # Horizontal logo just left of centre near top
    if LOGO_H.exists():
        logo_h = Emu(int(H * 0.06))
        logo_w = Emu(int(H * 0.06 * LOGO_H_RATIO))
        slide.shapes.add_picture(
            str(LOGO_H),
            Emu(int(W * 0.04)),
            Emu(int(H * 0.05)),
            logo_w, logo_h,
        )


# ── Section variants ──────────────────────────────────────────────────────

def _section(slide, prs, brand, **kw):
    W, H = prs.slide_width, prs.slide_height
    red = brand["colors"]["primary"]

    # Left vertical red bar, full height
    bar_w = Emu(int(W * 0.012))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, bar_w, H)
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(red); bar.line.fill.background()

    # Icon corner bottom-right — icon is taller than wide (W/H = 0.66)
    if LOGO_ICON.exists():
        icon_h = Emu(int(H * 0.08))
        icon_w = Emu(int(H * 0.08 * LOGO_ICON_RATIO))
        slide.shapes.add_picture(
            str(LOGO_ICON),
            W - icon_w - Emu(int(W * 0.04)),
            H - icon_h - Emu(int(H * 0.04)),
            icon_w, icon_h,
        )


def _section_filled(slide, prs, brand, **kw):
    """Full red background — for chapter-divider moments that need impact."""
    W, H = prs.slide_width, prs.slide_height
    red = brand["colors"]["primary"]

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = _rgb(red); bg.line.fill.background()
    # Push the background to the back so existing content sits on top
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)

    # White icon bottom-right — preserve aspect ratio
    if LOGO_ICON.exists():
        icon_h = Emu(int(H * 0.10))
        icon_w = Emu(int(H * 0.10 * LOGO_ICON_RATIO))
        slide.shapes.add_picture(
            str(LOGO_ICON),
            W - icon_w - Emu(int(W * 0.04)),
            H - icon_h - Emu(int(H * 0.04)),
            icon_w, icon_h,
        )


# ── Content variants ──────────────────────────────────────────────────────

def _content(slide, prs, brand, *,
             page_number=None, total_pages=None, section_label=None, **_):
    W, H = prs.slide_width, prs.slide_height
    red = brand["colors"]["primary"]

    # Bottom red accent bar — full width, thin
    bar_h = Emu(int(H * 0.012))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - bar_h, W, bar_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(red); bar.line.fill.background()

    # Page number — bottom-right, above the bar
    if page_number is not None:
        label = f"{page_number} / {total_pages}" if total_pages else str(page_number)
        _add_label(slide,
                   left=int(W * 0.92), top=int(H * 0.945),
                   width=int(W * 0.07), height=int(H * 0.04),
                   text=label, size_pt=9, color="#969696",
                   align=PP_ALIGN.RIGHT)

    # Small horizontal logo bottom-left
    if LOGO_H.exists():
        logo_h = Emu(int(H * 0.04))
        logo_w = Emu(int(H * 0.04 * LOGO_H_RATIO))
        slide.shapes.add_picture(
            str(LOGO_H),
            Emu(int(W * 0.04)),
            H - logo_h - Emu(int(H * 0.02)),
            logo_w, logo_h,
        )

    # Section eyebrow (top-right) — small label like "SOLUTION" or "WYTYCZNE"
    if section_label:
        _add_label(slide,
                   left=int(W * 0.75), top=int(H * 0.025),
                   width=int(W * 0.22), height=int(H * 0.04),
                   text=section_label.upper(), size_pt=9, bold=True,
                   color=brand["colors"]["primary"], align=PP_ALIGN.RIGHT)


def _content_minimal(slide, prs, brand, *,
                     page_number=None, total_pages=None, section_label=None, **_):
    """Minimalist chrome — thin top accent + page number only.

    For dense slides (charts, tables) where the full footer would compete
    with the data.
    """
    W, H = prs.slide_width, prs.slide_height
    red = brand["colors"]["primary"]

    # Top red accent — thin (4-px equivalent)
    bar_h = Emu(int(H * 0.006))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, bar_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(red); bar.line.fill.background()

    if page_number is not None:
        label = f"{page_number} / {total_pages}" if total_pages else str(page_number)
        _add_label(slide,
                   left=int(W * 0.92), top=int(H * 0.965),
                   width=int(W * 0.07), height=int(H * 0.03),
                   text=label, size_pt=9, color="#969696",
                   align=PP_ALIGN.RIGHT)

    if section_label:
        _add_label(slide,
                   left=int(W * 0.02), top=int(H * 0.965),
                   width=int(W * 0.50), height=int(H * 0.03),
                   text=section_label.upper(), size_pt=9, bold=True,
                   color=brand["colors"]["primary"], align=PP_ALIGN.LEFT)


def _closing(slide, prs, brand, **_):
    W, H = prs.slide_width, prs.slide_height
    red = brand["colors"]["primary"]

    # Thicker red bar bottom
    bar_h = Emu(int(H * 0.020))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - bar_h, W, bar_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(red); bar.line.fill.background()

    # Accent triangle bottom-left (mirror of cover_minimal)
    tri = Emu(int(W * 0.08))
    t = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, 0, H - bar_h - tri, tri, tri,
    )
    t.fill.solid(); t.fill.fore_color.rgb = _rgb(red)
    t.line.fill.background()
    t.rotation = 270

    # Larger horizontal logo bottom-right
    if LOGO_H.exists():
        logo_h = Emu(int(H * 0.07))
        logo_w = Emu(int(H * 0.07 * LOGO_H_RATIO))
        slide.shapes.add_picture(
            str(LOGO_H),
            W - logo_w - Emu(int(W * 0.05)),
            H - logo_h - Emu(int(H * 0.06)),
            logo_w, logo_h,
        )


# ── Helper ────────────────────────────────────────────────────────────────

def _add_label(slide, *, left, top, width, height,
               text, size_pt=10, color="#000000", bold=False,
               align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.name = "Raleway SemiBold" if bold else "Raleway"
        run.font.color.rgb = _rgb(color)


# ── Mode registry + intent mapping ────────────────────────────────────────

MODES = {
    "cover":            _cover_minimal,
    "cover_minimal":    _cover_minimal,
    "cover_split":      _cover_split,
    "cover_statement":  _cover_statement,
    "section":          _section,
    "section_filled":   _section_filled,
    "content":          _content,
    "content_minimal":  _content_minimal,
    "data_heavy":       _content_minimal,
    "closing":          _closing,
}


# Section / layout → default chrome mode
SECTION_MODES = {
    "cover":          "cover_statement",
    "contact":        "closing",
    "purpose":        "cover_minimal",
}

LAYOUT_MODES = {
    "cover":            "cover_statement",
    "section":          "section",
    "section_divider":  "section_filled",
    "closing":          "closing",
    "chart":            "content_minimal",
    "table_grid":       "content_minimal",
    "gantt":            "content_minimal",
    "comparison_matrix":"content_minimal",
    "big_quote":        "section_filled",
    "big_number":       "cover_minimal",
}


def mode_for(spec: dict) -> str:
    """Pick chrome mode; spec may override with `"chrome": "..."`."""
    if "chrome" in spec:
        return spec["chrome"]
    if "section" in spec and spec["section"] in SECTION_MODES:
        return SECTION_MODES[spec["section"]]
    layout = spec.get("layout") or spec.get("layout_override")
    if layout in LAYOUT_MODES:
        return LAYOUT_MODES[layout]
    return "content"
