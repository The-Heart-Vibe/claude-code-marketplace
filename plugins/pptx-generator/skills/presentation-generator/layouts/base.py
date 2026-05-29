"""
layouts/base.py — Helpers shared by every layout builder.

Centralises the boilerplate around python-pptx so per-layout files stay tiny:

* `add_slide(prs, layout_cfg)` — adds a slide using the right master layout
* `placeholder(slide, layout_cfg, role)` — fetches a shape by semantic role
* `set_text(shape, text, *, bullet=False)` — writes text while preserving
  the run-level font formatting that lives on the placeholder
* `clear(shape)` — empties a placeholder
* `set_bullets(shape, items)` — bullet list using the placeholder's prototype run
"""
from __future__ import annotations

from typing import Iterable, Optional

from pptx.enum.text import MSO_AUTO_SIZE


def add_slide(prs, layout_cfg: dict):
    """Add a slide using the Slide Master layout identified by master_idx."""
    idx = layout_cfg["master_idx"]
    layout = prs.slide_layouts[idx]
    return prs.slides.add_slide(layout)


def placeholder(slide, layout_cfg: dict, role: str):
    """Return the slide shape that corresponds to a semantic role, or None.

    Looks up the placeholder index from `layout_cfg['placeholders'][role]`
    and returns `slide.placeholders[idx]`. If the role isn't declared but
    a `roles_override` maps it to a known role, follow that mapping.
    Returns None if no placeholder can be resolved.
    """
    placeholders = layout_cfg.get("placeholders") or {}
    overrides = layout_cfg.get("roles_override") or {}

    if role in overrides:
        role = overrides[role]

    ph_idx = placeholders.get(role)
    if ph_idx is None:
        return None
    try:
        return slide.placeholders[ph_idx]
    except KeyError:
        return None


def _save_run_format(shape) -> dict:
    saved = {}
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name:
                saved["name"] = run.font.name
            if run.font.size:
                saved["size"] = run.font.size
            if run.font.bold is not None:
                saved["bold"] = run.font.bold
            try:
                if run.font.color.type:
                    saved["color"] = run.font.color.rgb
            except Exception:
                pass
            break
        if saved:
            break
    return saved


def _apply_run_format(run, saved: dict) -> None:
    """Re-apply saved run format, then force Raleway as the font face.

    The Slide Master often ships placeholders with Arial/Calibri as the
    placeholder prototype font. Saving + restoring keeps size/colour/bold
    but we override the font name unconditionally so every shape ends up
    in Raleway regardless of master defaults.
    """
    if "size" in (saved or {}):
        run.font.size = saved["size"]
    if "bold" in (saved or {}):
        run.font.bold = saved["bold"]
    if "color" in (saved or {}):
        try:
            run.font.color.rgb = saved["color"]
        except Exception:
            pass
    # Force brand font (overrides any saved name)
    run.font.name = "Raleway SemiBold" if run.font.bold else "Raleway"


def _enable_autofit(tf) -> None:
    """Apply word-wrap + shrink-to-fit to prevent text from bleeding past
    placeholder bounds when content is longer than the master layout
    anticipated."""
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def set_text(shape, text: Optional[str]) -> None:
    """Replace the placeholder's text while preserving font formatting.

    Accepts multi-line input (newlines become separate paragraphs).
    No-op if `shape` or `text` is falsy.
    """
    if shape is None or text is None:
        return
    text = str(text)
    fmt = _save_run_format(shape)
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            _apply_run_format(p.runs[0], fmt)
    _enable_autofit(tf)


def set_bullets(shape, items: Iterable[str]) -> None:
    """Render an iterable as bullet points in the placeholder."""
    if shape is None:
        return
    items = [i for i in (items or []) if i]
    if not items:
        return
    fmt = _save_run_format(shape)
    tf = shape.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        if p.runs:
            _apply_run_format(p.runs[0], fmt)
    _enable_autofit(tf)


def clear(shape) -> None:
    if shape is not None:
        shape.text_frame.clear()


def set_notes(slide, notes: Optional[str]) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
