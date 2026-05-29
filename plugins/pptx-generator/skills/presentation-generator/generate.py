"""
generate.py — Presentation generation engine (v2).

The Heart / VASBOX Plugin

Architecture (v2):
  * Loads a branded template .pptx whose Slide Master ships with named layouts.
  * Each spec entry resolves to a Slide Master layout (`prs.slide_layouts[idx]`)
    via `template_map.yaml`, NOT to a cloned slide.
  * Layout modules in `layouts/` fill placeholders by semantic role
    (`title`, `body`, `col1_heading`, ...) instead of by geometric heuristics.
  * Semantic pitch sections (`problem`, `solution`, ...) come from
    `narrative/pitchdeck.yaml` and pick the right generic layout for the
    content shape.

This removes ~500 lines of slide cloning, position heuristics, and
example-text scrubbing that v1 needed because it cloned full template slides.
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Callable, Optional

import yaml
from pptx import Presentation

import layouts as layouts_pkg

# ── Paths ──────────────────────────────────────────────────────────────────
BUNDLE_DIR    = Path(__file__).parent
ICONS_DIR     = BUNDLE_DIR / "icons"
SVG_DIR       = ICONS_DIR / "svgs"
CACHE_DIR     = ICONS_DIR / "cache"
BRAND_FILE    = BUNDLE_DIR / "brand.yaml"
TMAP_FILE     = BUNDLE_DIR / "template_map.yaml"
NARRATIVE_DIR = BUNDLE_DIR / "narrative"
TEMPLATES_DIR = BUNDLE_DIR / "templates"

DEFAULT_TEMPLATE = "pitchdeck-toolkit"


# ── Config loaders ─────────────────────────────────────────────────────────

def load_brand() -> dict:
    with open(BRAND_FILE) as f:
        return yaml.safe_load(f)


def load_template_map() -> dict:
    if not TMAP_FILE.exists():
        return {"templates": {}}
    with open(TMAP_FILE) as f:
        return yaml.safe_load(f) or {"templates": {}}


def load_narrative(name: str = "pitchdeck") -> dict:
    path = NARRATIVE_DIR / f"{name}.yaml"
    if not path.exists():
        return {}
    with open(path) as f:
        return yaml.safe_load(f) or {}


def resolve_template_path(name: str) -> Path:
    p = Path(name)
    if p.suffix == ".pptx" and p.is_absolute() and p.exists():
        return p
    stem = p.stem if p.suffix == ".pptx" else name
    candidate = TEMPLATES_DIR / f"{stem}.pptx"
    if candidate.exists():
        return candidate
    raise FileNotFoundError(
        f"Template not found: {name!r} (looked at {candidate})"
    )


# ── Layout cfg resolver ────────────────────────────────────────────────────

def _resolve_layout_cfg(tmap: dict, tname: str, alias: str) -> dict:
    layouts = (
        tmap.get("templates", {})
            .get(tname, {})
            .get("layouts", {})
    )
    cfg = layouts.get(alias)
    if cfg is None:
        raise KeyError(
            f"Layout alias {alias!r} not defined for template {tname!r}. "
            f"Edit template_map.yaml."
        )
    # follow alias chain (e.g. cover -> cover_split)
    if "alias_of" in cfg:
        return _resolve_layout_cfg(tmap, tname, cfg["alias_of"])
    return cfg


# ── Icon engine ────────────────────────────────────────────────────────────

def _load_icon_manifest() -> dict:
    with open(ICONS_DIR / "icon_manifest.json") as f:
        return json.load(f)


def _resolve_icon_name(keyword: str, manifest: dict) -> str:
    kw = keyword.lower().strip()
    aliases = manifest.get("aliases", {})
    if kw in aliases:
        return aliases[kw]
    if (SVG_DIR / f"{kw}.svg").exists():
        return kw
    for alias, icon in aliases.items():
        if kw in alias or alias in kw:
            return icon
    return "star"


def get_icon_png(keyword: str, size: int = 64, color: str = "#1A1A2E") -> Optional[Path]:
    try:
        import cairosvg
    except ImportError:
        print("  [icon] cairosvg not installed — skipping icons")
        return None
    CACHE_DIR.mkdir(exist_ok=True)
    manifest = _load_icon_manifest()
    name = _resolve_icon_name(keyword, manifest)
    svg = SVG_DIR / f"{name}.svg"
    if not svg.exists():
        return None
    cache = CACHE_DIR / f"{name}_{size}_{color.lstrip('#')}.png"
    if not cache.exists():
        text = svg.read_text()
        text = text.replace('stroke="currentColor"', f'stroke="{color}"')
        text = text.replace("stroke:currentColor",  f"stroke:{color}")
        cairosvg.svg2png(
            bytestring=text.encode(),
            write_to=str(cache),
            output_width=size,
            output_height=size,
        )
    return cache


# ── Chart engine ───────────────────────────────────────────────────────────

def render_chart(chart_type: str, data: dict, brand: dict, output_path: str) -> Optional[str]:
    """Render a chart to PNG. Supported types:
       bar, line, pie, donut, horizontal_bar, stacked_bar, waterfall, funnel.
    """
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np
    except ImportError:
        return None

    # Force Raleway font on matplotlib (falls back to sans-serif if absent)
    body_font = brand.get("typography", {}).get("body_font", "Raleway")
    plt.rcParams["font.family"] = [body_font, "Arial", "sans-serif"]
    plt.rcParams["font.size"] = 11

    labels  = data.get("labels", [])
    values  = data.get("values", [])
    series  = data.get("series", [])   # for stacked_bar
    primary = brand["colors"]["primary"]
    bcolors = brand["colors"]
    palette = [
        bcolors.get("primary",   "#E61B25"),
        bcolors.get("blue",      "#0056A4"),
        bcolors.get("green",     "#13A538"),
        bcolors.get("red_light", "#E9787E"),
        bcolors.get("gray_1",    "#969696"),
        bcolors.get("black",     "#000000"),
    ]

    fig, ax = plt.subplots(figsize=(10, 5.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("#FAFAFA")
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#E5E7EB")
    ax.tick_params(colors="#6B7280")

    if chart_type == "bar":
        ax.bar(labels, values, color=primary, zorder=3)
        ax.yaxis.grid(True, color="#E5E7EB", zorder=0)
        if values:
            ax.set_ylim(0, max(values) * 1.2)
        for i, v in enumerate(values):
            ax.text(i, v, f"{v}", ha="center", va="bottom",
                    color=primary, fontsize=10, fontweight="bold")

    elif chart_type == "horizontal_bar":
        y_pos = np.arange(len(labels))
        ax.barh(y_pos, values, color=primary, zorder=3)
        ax.set_yticks(y_pos)
        ax.set_yticklabels(labels)
        ax.invert_yaxis()
        ax.xaxis.grid(True, color="#E5E7EB", zorder=0)

    elif chart_type == "stacked_bar":
        # series is list of {label, values}
        bottom = np.zeros(len(labels))
        for i, s in enumerate(series):
            ax.bar(labels, s["values"], bottom=bottom,
                   label=s.get("label", f"Series {i+1}"),
                   color=palette[i % len(palette)], zorder=3)
            bottom += np.array(s["values"])
        ax.yaxis.grid(True, color="#E5E7EB", zorder=0)
        ax.legend(loc="upper left", frameon=False)

    elif chart_type == "line":
        if series:
            for i, s in enumerate(series):
                ax.plot(labels, s["values"], linewidth=2.5,
                        marker="o", color=palette[i % len(palette)],
                        label=s.get("label", f"Series {i+1}"))
            ax.legend(loc="upper left", frameon=False)
        else:
            ax.plot(labels, values, color=primary, linewidth=2.5,
                    marker="o", markerfacecolor=primary, zorder=3)
        ax.yaxis.grid(True, color="#E5E7EB", zorder=0)

    elif chart_type == "pie":
        ax.pie(values, labels=labels, colors=palette[:len(values)],
               autopct="%1.0f%%", startangle=90,
               wedgeprops={"linewidth": 2, "edgecolor": "white"})
        ax.axis("equal")

    elif chart_type == "donut":
        wedges, _ = ax.pie(values, labels=labels, colors=palette[:len(values)],
                           startangle=90,
                           wedgeprops={"linewidth": 2, "edgecolor": "white",
                                       "width": 0.35})
        ax.axis("equal")
        # Centre stat
        total = sum(values)
        ax.text(0, 0, f"{total}",
                ha="center", va="center",
                fontsize=28, fontweight="bold", color=primary)

    elif chart_type == "waterfall":
        # values can include negatives; first and last are "totals"
        cum = 0
        bottoms = []
        bar_colors = []
        for i, v in enumerate(values):
            if i == 0 or i == len(values) - 1:
                bottoms.append(0)
                bar_colors.append(bcolors.get("gray_1", "#969696"))
                cum = v if i == 0 else cum
            else:
                bottoms.append(cum)
                bar_colors.append(
                    bcolors.get("green", "#13A538") if v >= 0
                    else bcolors.get("primary", "#E61B25")
                )
                cum += v
        ax.bar(labels, [abs(v) for v in values], bottom=bottoms,
               color=bar_colors, zorder=3)
        ax.yaxis.grid(True, color="#E5E7EB", zorder=0)

    elif chart_type == "funnel":
        # Each value drawn as a horizontal bar, widest first
        sorted_pairs = sorted(zip(labels, values), key=lambda x: -x[1])
        ys = list(range(len(sorted_pairs)))
        ys = ys[::-1]
        for (lab, val), y in zip(sorted_pairs, ys):
            ax.barh(y, val, color=palette[(len(ys) - y - 1) % len(palette)],
                    zorder=3)
            ax.text(val / 2, y, f"{lab}: {val}",
                    ha="center", va="center", color="white",
                    fontsize=11, fontweight="bold")
        ax.set_yticks([])
        ax.set_xticks([])
        ax.spines[["left", "bottom"]].set_visible(False)

    else:
        ax.text(0.5, 0.5, f"Unsupported chart type: {chart_type!r}",
                ha="center", va="center", transform=ax.transAxes)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches="tight",
                facecolor="white")
    plt.close()
    return output_path


# ── Main generation ───────────────────────────────────────────────────────

def generate(spec: dict, output_path: str, template_name: Optional[str] = None) -> str:
    brand = load_brand()
    tmap = load_template_map()
    tname = template_name or spec.get("template") or DEFAULT_TEMPLATE
    tpath = resolve_template_path(tname)
    output_dir = Path(output_path).parent

    print(f"  [template] {tpath.name}")

    # Push brand typography into the widget layer so every free text box
    # gets the right font (Raleway), instead of falling back to Arial.
    from layouts import widgets
    typo = brand.get("typography", {})
    widgets.set_default_fonts(
        body=typo.get("body_font"),
        emphasis=typo.get("emphasis_font") or typo.get("heading_font"),
        light=typo.get("light_font"),
        fallback=typo.get("fallback_font"),
    )

    # Open the template; we add fresh slides from its Slide Master layouts.
    prs = Presentation(str(tpath))
    n_original = len(prs.slides)

    # Narrative selection: blank.pptx → general (internal decks);
    # pitchdeck-toolkit → pitchdeck (investor decks). A spec may override
    # via `"narrative": "<name>"`.
    narrative_name = spec.get("narrative") or _default_narrative_for(tname)
    narrative = load_narrative(narrative_name)

    def resolve_layout(alias: str) -> dict:
        return _resolve_layout_cfg(tmap, tname, alias)

    ctx = {
        "get_icon_png":    get_icon_png,
        "render_chart":    render_chart,
        "output_dir":      output_dir,
        "narrative":       narrative,
        "resolve_layout":  resolve_layout,
    }

    from layouts import brand_chrome

    slides_specs = spec.get("slides", [])
    total = len(slides_specs)

    for i, slide_spec in enumerate(slides_specs, start=1):
        layout_alias = _layout_alias_for(slide_spec, narrative)
        layout_cfg = resolve_layout(layout_alias)
        builder: Callable = layouts_pkg.resolve(slide_spec)
        label = (
            slide_spec.get("title")
            or slide_spec.get("section")
            or slide_spec.get("layout", "")
        )
        print(f"  building: {layout_alias:<14} {label[:55]}")
        slide = builder.build(prs, layout_cfg, slide_spec, brand, ctx)

        # Clear placeholders the builder left empty so the master's prompt
        # text ("Click to add text") doesn't render in the final deck.
        _clear_empty_placeholders(slide)

        # Apply brand chrome (red bar, logo, page numbers, section eyebrow)
        # so output decks remain on-brand even though templates were stripped.
        mode = brand_chrome.mode_for(slide_spec)
        is_cover_or_section = mode.startswith("cover") or mode.startswith("section")
        page_number = None if is_cover_or_section else i

        # Derive a section label for the top-right eyebrow on content slides.
        section_name = slide_spec.get("section")
        section_label = None
        if section_name and not is_cover_or_section:
            section = (narrative.get("sections") or {}).get(section_name, {})
            section_label = section.get("title_en") or section_name.replace("_", " ")

        brand_chrome.apply(
            slide, prs, brand,
            mode=mode,
            page_number=page_number,
            total_pages=total,
            section_label=section_label,
        )

    # Drop the original template slides we inherited from the .pptx
    _purge_original_slides(prs, n_original)

    prs.save(output_path)
    print(f"\n  done → {output_path}  ({len(prs.slides)} slides)")
    return output_path


def _default_narrative_for(template_name: str) -> str:
    """Map a template to its default narrative YAML."""
    return {
        "pitchdeck-toolkit": "pitchdeck",
        "blank":             "general",
    }.get(template_name, "pitchdeck")


def _clear_empty_placeholders(slide) -> None:
    """Set empty placeholders to a non-prompt state.

    PowerPoint renders the master's prompt text ("Click to add text") on any
    placeholder that's still empty. We can't delete the placeholder safely
    (it lives in the layout), but we can mark the prompt as unwanted by
    removing the placeholder element's <p:ph> hint so the renderer treats it
    as a normal empty shape.
    """
    from pptx.oxml.ns import qn
    for ph in list(slide.placeholders):
        if not ph.has_text_frame:
            continue
        if ph.text_frame.text.strip():
            continue
        # Remove the inheritance hint so the master prompt doesn't render.
        nvSpPr = ph._element.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find(qn("p:nvPr"))
        if nvPr is None:
            continue
        ph_hint = nvPr.find(qn("p:ph"))
        if ph_hint is not None:
            nvPr.remove(ph_hint)


def _layout_alias_for(spec: dict, narrative: dict) -> str:
    """Pick the alias that template_map.yaml should resolve for this spec."""
    if "layout_override" in spec:
        return spec["layout_override"]
    if "section" in spec:
        sec = narrative.get("sections", {}).get(spec["section"], {})
        return sec.get("preferred_layout", "content")
    return spec.get("layout", "content")


def _purge_original_slides(prs, n_original: int) -> None:
    """Remove the slides that were already in the template when we opened it.

    Removes from the back so earlier indices stay stable.
    """
    from pptx.oxml.ns import qn
    sld_id_lst = prs.slides._sldIdLst
    slides = list(prs.slides)
    for i in range(n_original - 1, -1, -1):
        slide = slides[i]
        rId = None
        for rel_id, rel in prs.part.rels.items():
            if not rel.is_external and rel.target_part is slide.part:
                rId = rel_id
                break
        if rId is None:
            continue
        for sld_id in list(sld_id_lst):
            if sld_id.get(qn("r:id")) == rId:
                sld_id_lst.remove(sld_id)
                break
        prs.part.drop_rel(rId)


# ── CLI ────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate a branded presentation")
    parser.add_argument("--spec",     required=True, help="Path to spec JSON")
    parser.add_argument("--output",   required=True, help="Output .pptx path")
    parser.add_argument(
        "--template",
        default=None,
        help=f"Template name. Default: {DEFAULT_TEMPLATE}",
    )
    args = parser.parse_args()

    with open(args.spec) as f:
        spec = json.load(f)

    generate(spec, args.output, template_name=args.template)
